import pandas as pd
import streamlit as st
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches
from docx.shared import Pt
import matplotlib.pyplot as plt


# Reports Page
def reports_page():
    st.image('bkg.png', use_column_width=True)
    st.title('Threat Analysis')
    st.write('This page displays KPI metrics and charts based on the Network Anomaly Dataset.')

    # Load dataset (you might need to adjust the path)
    df = pd.read_csv('Network_anomaly_data.csv')

    # Filters
    st.sidebar.title('Filters')
    attack_type_filter = st.sidebar.multiselect('Select Attack Types', df['attack_type'].unique(),
                                                default=df['attack_type'].unique())
    protocol_type_filter = st.sidebar.multiselect('Select Protocol Types', df['protocoltype'].unique(),
                                                  default=df['protocoltype'].unique())

    # Apply filters to the dataframe
    data = df[(df['attack_type'].isin(attack_type_filter)) & (df['protocoltype'].isin(protocol_type_filter))]

    # Add spacing
    st.markdown("---")

    # KPI Metrics
    st.header('KPI Metrics')

    # Example KPI metrics
    num_records = len(data)
    num_attacks = len(data[data['attack'] != 'normal'])
    attack_percentage = (num_attacks / num_records) * 100

    col1, col2, col3 = st.columns(3)
    with col1:
        st.metric("Total Connections", f"|{num_records}|")
    with col2:
        st.metric("Total Attacks", f"|{num_attacks}|")
    with col3:
        st.metric("Percentage of Attacks", f"|{attack_percentage:.2f}%|")

    # Add spacing
    st.markdown("---")

    # Example Charts
    st.header('Charts')
    # Attack types distribution
    attack_types = data['attack_type'].value_counts()
    st.subheader('Distribution of Attack Types')
    st.bar_chart(attack_types)

    # Plotting with matplotlib
    st.subheader('Distribution of Protocol Types')
    protocol_counts = data['protocoltype'].value_counts()
    st.bar_chart(protocol_counts)

    # Export filtered data
    st.sidebar.markdown("## Export Data")
    csv = data.to_csv(index=False)
    st.sidebar.download_button(
        label="⬇ Download CSV File",
        data=csv,
        file_name='filtered_data.csv',
        mime='text/csv',
    )

    # Create Matplotlib charts
    def create_charts():
        loc_fig, ax = plt.subplots(1, 2, figsize=(14, 6))

        # Attack types distribution
        attack_types.plot(kind='bar', ax=ax[0], color='skyblue')
        ax[0].set_title('Distribution of Attack Types')
        ax[0].set_xlabel('Attack Type')
        ax[0].set_ylabel('Count')

        # Protocol types distribution
        protocol_counts.plot(kind='bar', ax=ax[1], color='lightgreen')
        ax[1].set_title('Distribution of Protocol Types')
        ax[1].set_xlabel('Protocol Type')
        ax[1].set_ylabel('Count')

        plt.tight_layout()
        return loc_fig

    fig = create_charts()

    # Save the charts as an image
    img_buf = BytesIO()
    fig.savefig(img_buf, format='png')
    img_buf.seek(0)

    # st.image(img_buf, caption='Charts', use_column_width=True)

    st.sidebar.markdown("## Export Report")

    # Export to PowerPoint
    def create_ppt(imgbuf):
        prs = Presentation()
        slide_layout = prs.slide_layouts[5]  # Title and Content

        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Network Anomaly Report"

        # Add KPI Metrics
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
        text_frame = textbox.text_frame
        p = text_frame.add_paragraph()
        p.text = f'Total Connections: {num_records}'
        p.font.size = Pt(18)
        p = text_frame.add_paragraph()
        p.text = f'Total Attacks: {num_attacks}'
        p.font.size = Pt(18)
        p = text_frame.add_paragraph()
        p.text = f'Percentage of Attacks: {attack_percentage:.2f}%'
        p.font.size = Pt(18)

        # Add Chart
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Charts"
        slide.shapes.add_picture(imgbuf, Inches(1), Inches(1.5), width=Inches(8), height=Inches(4.5))

        return prs

    pptx = create_ppt(img_buf)
    buffer = BytesIO()
    pptx.save(buffer)
    buffer.seek(0)

    st.sidebar.download_button(
        label="⬇ Download PowerPoint",
        data=buffer,
        file_name='Network_Anomaly_Report.pptx',
        mime='application/vnd.openxmlformats-officedocument.presentationml.presentation',
    )
